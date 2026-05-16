"""Example conference-deck PPTX generator.

Reference implementation showing the slide patterns the conference-deck skill
expects: title wall, three-column framing, bullet-led case study slides,
horizontal flow, KPI grid, pull quote with accent bar, pattern-comparison
table, and a numbered take-home slide.

Copy this file, rename it to `build_[slug]_deck.py` in the `scripts/`
directory, and replace placeholder content with your actual deck content.

Requires: pip install python-pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path(__file__).resolve().parents[3] / "output" / "decks" / "example-conference-deck.pptx"

# Palette — adjust per brand_guidelines.md
BG = RGBColor(0x0A, 0x0F, 0x1C)       # deep background
INK = RGBColor(0xF5, 0xF7, 0xFA)      # primary text
MUTED = RGBColor(0x9A, 0xA3, 0xB2)    # secondary text
ACCENT = RGBColor(0x3B, 0xE8, 0xB0)   # accent 1 (case study 1 color)
ACCENT2 = RGBColor(0xFF, 0x6A, 0x3D)  # accent 2 (case study 2 color)
ROW_ALT = RGBColor(0x13, 0x1A, 0x2B)  # tile / alt row fill

FONT = "Helvetica Neue"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------- Helper functions ----------

def paint_background(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_textbox(slide, left, top, width, height, text, size=24, bold=False,
                color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, size=22, color=INK,
                bullet_char="\u2022", spacing_before=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(spacing_before if i > 0 else 0)
        run = p.add_run()
        run.text = f"{bullet_char}  {item}" if bullet_char else item
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6), color=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar


def add_section_label(slide, text, left=Inches(0.6), top=Inches(0.5), color=ACCENT):
    add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.35), color=color)
    add_textbox(slide, left + Inches(0.2), top - Inches(0.02), Inches(8), Inches(0.4),
                text.upper(), size=12, bold=True, color=color)


def add_slide_title(slide, text, top=Inches(1.05)):
    add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.9),
                text, size=36, bold=True, color=INK)


def add_footer(slide, text="[Event name]  ·  [Year or hashtag]"):
    add_textbox(slide, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
                text, size=10, color=MUTED)


def set_notes(slide, notes_text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = notes_text


def add_table(slide, left, top, width, height, data, header=True,
              header_fill=ACCENT, header_text=BG, body_text=INK,
              body_fill=BG, alt_fill=ROW_ALT, font_size=14):
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.fill.solid()
            if header and r == 0:
                cell.fill.fore_color.rgb = header_fill
                text_color = header_text
                bold = True
            else:
                cell.fill.fore_color.rgb = alt_fill if (r % 2 == 0) else body_fill
                text_color = body_text
                bold = False
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.12)
            tf.margin_right = Inches(0.12)
            tf.margin_top = Inches(0.08)
            tf.margin_bottom = Inches(0.08)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.name = FONT
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = text_color
    return table


def add_kpi_grid(slide, left, top, cell_w, cell_h, items, cols=3, gap=Inches(0.2)):
    """items: list of (label, value) tuples. Value rendered large in accent; label small in muted."""
    for i, (label, value) in enumerate(items):
        r = i // cols
        c = i % cols
        x = left + c * (cell_w + gap)
        y = top + r * (cell_h + gap)
        tile = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cell_w, cell_h)
        tile.line.color.rgb = MUTED
        tile.line.width = Pt(0.5)
        tile.fill.solid()
        tile.fill.fore_color.rgb = ROW_ALT
        add_textbox(slide, x + Inches(0.15), y + Inches(0.15), cell_w - Inches(0.3), Inches(0.8),
                    str(value), size=28, bold=True, color=ACCENT)
        add_textbox(slide, x + Inches(0.15), y + Inches(1.05), cell_w - Inches(0.3), cell_h - Inches(1.1),
                    str(label), size=12, color=MUTED)


def add_flow(slide, left, top, steps, box_w=Inches(2.7), box_h=Inches(0.7), gap=Inches(0.3)):
    """Horizontal flow: boxes connected by right arrows."""
    for i, step in enumerate(steps):
        x = left + i * (box_w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        box.line.color.rgb = ACCENT
        box.fill.solid()
        box.fill.fore_color.rgb = BG
        tf = box.text_frame
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step
        run.font.name = FONT
        run.font.size = Pt(14)
        run.font.color.rgb = INK
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                            x + box_w + Inches(0.02), top + Inches(0.27),
                                            Inches(0.26), Inches(0.16))
            arrow.line.fill.background()
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT


# ---------- Build deck ----------

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# Slide 1: Title
s = prs.slides.add_slide(blank)
paint_background(s)
add_accent_bar(s, Inches(0.6), Inches(2.6), width=Inches(1.2), height=Inches(0.08), color=ACCENT)
add_textbox(s, Inches(0.6), Inches(2.8), Inches(12), Inches(1.4),
            "[Session Title]", size=64, bold=True, color=INK)
add_textbox(s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.8),
            "[Subtitle or thesis]", size=28, color=MUTED)
add_textbox(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
            "[Presenter name]  ·  [Title]  ·  [Company]", size=16, color=INK)
add_footer(s)
set_notes(s, "Opening line: [one-sentence premise]. Introduce yourself briefly.")


# Slide 2: Premise (three short lines)
s = prs.slides.add_slide(blank)
paint_background(s)
add_section_label(s, "THE PREMISE")
add_slide_title(s, "[The shift is already here.]")
items = ["[Claim 1.]", "[Claim 2.]", "[Claim 3.]"]
top = Inches(2.6)
for i, line in enumerate(items):
    add_textbox(s, Inches(0.6), top + i * Inches(1.0), Inches(12), Inches(0.9),
                line, size=34, bold=True, color=INK)
add_footer(s)
set_notes(s, "Anchor the audience. State what is different about right now. Keep under 45 seconds.")


# Slide 3: Framing (three-column questions or themes)
s = prs.slides.add_slide(blank)
paint_background(s)
add_section_label(s, "FRAMING")
add_slide_title(s, "[What this talk will answer.]")
questions = [
    ("[THEME 1]", "[Question or statement 1?]"),
    ("[THEME 2]", "[Question or statement 2?]"),
    ("[THEME 3]", "[Question or statement 3?]"),
]
col_w = Inches(4.0)
gap = Inches(0.35)
start_x = Inches(0.6)
start_y = Inches(2.6)
for i, (label, q) in enumerate(questions):
    x = start_x + i * (col_w + gap)
    tile = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y, col_w, Inches(3.6))
    tile.line.color.rgb = MUTED
    tile.line.width = Pt(0.5)
    tile.fill.solid()
    tile.fill.fore_color.rgb = ROW_ALT
    add_textbox(s, x + Inches(0.3), start_y + Inches(0.2), col_w - Inches(0.6), Inches(0.6),
                f"0{i+1}", size=36, bold=True, color=ACCENT)
    add_textbox(s, x + Inches(0.3), start_y + Inches(0.9), col_w - Inches(0.6), Inches(0.5),
                label, size=14, bold=True, color=MUTED)
    add_textbox(s, x + Inches(0.3), start_y + Inches(1.5), col_w - Inches(0.6), Inches(2.0),
                q, size=18, color=INK)
add_footer(s)
set_notes(s, "Map each theme to a case study so the audience knows the structure ahead.")


# Slide 4: Case study 1 — problem
s = prs.slides.add_slide(blank)
paint_background(s)
add_section_label(s, "CASE STUDY  ·  [COMPANY 1]  ·  [THEME 1]", color=ACCENT)
add_slide_title(s, "[The problem they faced.]")
add_bullets(s, Inches(0.6), Inches(2.5), Inches(12), Inches(4), [
    "[Specific problem 1, with context.]",
    "[Specific problem 2.]",
    "[Specific problem 3.]",
    "[Compliance or stakes pressure.]",
], size=22, spacing_before=14)
add_footer(s)
set_notes(s, "Summarize the problem in the customer's own framing. Cite them if you can.")


# Slide 5: Case study 1 — deployment + flow
s = prs.slides.add_slide(blank)
paint_background(s)
add_section_label(s, "CASE STUDY  ·  [COMPANY 1]  ·  [THEME 1]", color=ACCENT)
add_slide_title(s, "What they deployed.")
add_bullets(s, Inches(0.6), Inches(2.5), Inches(12), Inches(3), [
    "[Capability 1.]",
    "[Capability 2.]",
    "[Capability 3.]",
    "[Capability 4.]",
], size=22, spacing_before=14)
add_flow(s, Inches(0.6), Inches(5.7), ["[Step 1]", "[Step 2]", "[Step 3]", "[Step 4]"])
add_footer(s)
set_notes(s, "Walk through the deployment in four beats. The flow at the bottom visualizes the sequence.")


# Slide 6: Case study 1 — outcomes (KPI grid)
s = prs.slides.add_slide(blank)
paint_background(s)
add_section_label(s, "CASE STUDY  ·  [COMPANY 1]  ·  [THEME 1]", color=ACCENT)
add_slide_title(s, "What changed.")
add_kpi_grid(s, Inches(0.6), Inches(2.4), Inches(4.0), Inches(1.9), [
    ("[Label 1]", "[Value 1]"),
    ("[Label 2]", "[Value 2]"),
    ("[Label 3]", "[Value 3]"),
    ("[Label 4]", "[Value 4]"),
    ("[Label 5]", "[Value 5]"),
    ("[Label 6]", "[Value 6]"),
], cols=3)
add_footer(s)
set_notes(s, "Land the most memorable number first. Every value here must be verbatim from the case study.")


# Slide 7: Case study 1 — pull quote
s = prs.slides.add_slide(blank)
paint_background(s)
add_section_label(s, "CASE STUDY  ·  [COMPANY 1]  ·  [THEME 1]", color=ACCENT)
add_accent_bar(s, Inches(0.6), Inches(2.4), width=Inches(0.12), height=Inches(3.2), color=ACCENT)
add_textbox(s, Inches(1.0), Inches(2.4), Inches(11.5), Inches(3.0),
            "\u201C[Verbatim quote from a named customer employee.]\u201D",
            size=30, color=INK)
add_textbox(s, Inches(1.0), Inches(6.0), Inches(11.5), Inches(0.4),
            "[Name]  ·  [Title]  ·  [Company]",
            size=14, color=MUTED)
add_footer(s)
set_notes(s, "Let the quote breathe. Transition to case study 2 in the notes, not on the slide.")


# Slide 8: Case study 2 — problem (different accent color)
s = prs.slides.add_slide(blank)
paint_background(s)
add_section_label(s, "CASE STUDY  ·  [COMPANY 2]  ·  [THEME 2]", color=ACCENT2)
add_slide_title(s, "[The problem they faced.]")
add_bullets(s, Inches(0.6), Inches(2.5), Inches(12), Inches(4), [
    "[Problem 1.]",
    "[Problem 2.]",
    "[Problem 3.]",
    "[Problem 4.]",
], size=22, spacing_before=14)
add_footer(s)
set_notes(s, "Set the contrast with case study 1. Different industry or angle.")


# Slide 9: Case study 2 — deployment
s = prs.slides.add_slide(blank)
paint_background(s)
add_section_label(s, "CASE STUDY  ·  [COMPANY 2]  ·  [THEME 2]", color=ACCENT2)
add_slide_title(s, "What they deployed.")
add_bullets(s, Inches(0.6), Inches(2.5), Inches(12), Inches(4), [
    "[Capability 1.]",
    "[Capability 2.]",
    "[Capability 3.]",
    "[Capability 4.]",
], size=22, spacing_before=14)
add_footer(s)
set_notes(s, "Emphasize the unique angle this customer took. Do not repeat slide 5.")


# Slide 10: Case study 2 — outcomes
s = prs.slides.add_slide(blank)
paint_background(s)
add_section_label(s, "CASE STUDY  ·  [COMPANY 2]  ·  [THEME 2]", color=ACCENT2)
add_slide_title(s, "What changed.")
add_kpi_grid(s, Inches(0.6), Inches(2.4), Inches(4.0), Inches(1.9), [
    ("[Label 1]", "[Value 1]"),
    ("[Label 2]", "[Value 2]"),
    ("[Label 3]", "[Value 3]"),
    ("[Label 4]", "[Value 4]"),
    ("[Label 5]", "[Value 5]"),
    ("[Label 6]", "[Value 6]"),
], cols=3)
add_footer(s)
set_notes(s, "Pick one number to dwell on. One sentence landing.")


# Slide 11: Case study 2 — pull quote
s = prs.slides.add_slide(blank)
paint_background(s)
add_section_label(s, "CASE STUDY  ·  [COMPANY 2]  ·  [THEME 2]", color=ACCENT2)
add_accent_bar(s, Inches(0.6), Inches(2.4), width=Inches(0.12), height=Inches(3.2), color=ACCENT2)
add_textbox(s, Inches(1.0), Inches(2.4), Inches(11.5), Inches(3.0),
            "\u201C[Verbatim quote from a named employee.]\u201D",
            size=28, color=INK)
add_textbox(s, Inches(1.0), Inches(6.0), Inches(11.5), Inches(0.4),
            "[Name]  ·  [Title]  ·  [Company]",
            size=14, color=MUTED)
add_footer(s)
set_notes(s, "Transition into the synthesis: these two companies did not coordinate, they landed on the same architecture.")


# Slide 12: Pattern / synthesis table
s = prs.slides.add_slide(blank)
paint_background(s)
add_section_label(s, "THE PATTERN")
add_slide_title(s, "[The pattern across both cases.]")
data = [
    ["Principle", "[Company 1] ([Theme 1])", "[Company 2] ([Theme 2])"],
    ["[Principle 1]", "[How it showed up]", "[How it showed up]"],
    ["[Principle 2]", "[How it showed up]", "[How it showed up]"],
    ["[Principle 3]", "[How it showed up]", "[How it showed up]"],
    ["[Principle 4]", "[How it showed up]", "[How it showed up]"],
]
add_table(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(4.0), data, font_size=15)
add_footer(s)
set_notes(s, "This is the takeaway slide in disguise. The audience should be able to photograph this and walk out with the architecture.")


# Slide 13: Big idea / accountability
s = prs.slides.add_slide(blank)
paint_background(s)
add_section_label(s, "[THE BIG IDEA]")
add_accent_bar(s, Inches(0.6), Inches(2.1), width=Inches(1.4), height=Inches(0.1), color=ACCENT)
add_textbox(s, Inches(0.6), Inches(2.3), Inches(12), Inches(1.8),
            "[Punchline line 1].\n[Punchline line 2].",
            size=64, bold=True, color=INK)
add_bullets(s, Inches(0.6), Inches(5.4), Inches(12), Inches(1.6), [
    "[Supporting point 1.]",
    "[Supporting point 2.]",
    "[Supporting point 3.]",
], size=16, color=MUTED, spacing_before=10)
add_footer(s)
set_notes(s, "Answer the abstract's hardest question here. Do not soften.")


# Slide 14: Take home + CTA
s = prs.slides.add_slide(blank)
paint_background(s)
add_section_label(s, "TAKE HOME")
add_slide_title(s, "[N] things to walk out with.")
takeaways = [
    "[Takeaway 1 mapped to abstract bullet 1.]",
    "[Takeaway 2 mapped to abstract bullet 2.]",
    "[Takeaway 3.]",
    "[Takeaway 4.]",
    "[Takeaway 5.]",
]
top = Inches(2.4)
for i, line in enumerate(takeaways):
    num_box = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), top + i * Inches(0.85),
                                  Inches(0.55), Inches(0.55))
    num_box.line.fill.background()
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = ACCENT
    tf = num_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(i + 1)
    run.font.name = FONT
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = BG
    add_textbox(s, Inches(1.3), top + i * Inches(0.85) + Inches(0.05),
                Inches(11.5), Inches(0.75), line, size=18, color=INK)
add_textbox(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
            "[Presenter name]  ·  [Company]  ·  [URL]",
            size=12, color=MUTED)
add_footer(s)
set_notes(s, "Read these slowly. Close with an invitation: booth, hallway, office hours.")


OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
